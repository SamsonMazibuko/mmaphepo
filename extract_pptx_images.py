import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_FILE = "Mmaphepo Holdings Product Bookltet.pptx"
OUTPUT_DIR = "mmaphepo-website/public/images"

os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation(PPTX_FILE)
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            image_format = image.ext
            image_filename = f"slide_{i+1}_img_{shape.shape_id}.{image_format}"
            image_path = os.path.join(OUTPUT_DIR, image_filename)
            with open(image_path, 'wb') as img_file:
                img_file.write(image_bytes)
print(f"Images extracted to {OUTPUT_DIR}/") 