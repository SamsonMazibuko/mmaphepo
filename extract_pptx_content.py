import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_FILE = "Mmaphepo Holdings Product Bookltet.pptx"
OUTPUT_TEXT = "pptx_text_output.txt"

prs = Presentation(PPTX_FILE)
with open(OUTPUT_TEXT, 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(shape.text.strip() + "\n")
print(f"Text extracted to {OUTPUT_TEXT}") 